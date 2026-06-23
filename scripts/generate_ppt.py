import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Configure 16:9 widescreen layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ---------------------------------------------------------
    # SLIDE 1: Dark Pine Green Theme
    # ---------------------------------------------------------
    blank_slide_layout = prs.slide_layouts[6] # completely blank layout
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Background color
    background1 = slide1.background
    fill1 = background1.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(19, 48, 43) # #13302B (Pine Green)
    
    # TOP HEADER TEXT (MINI PROJECT ABSTRACT PRESENTATION)
    header_box = slide1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8), Inches(0.5))
    tf_h = header_box.text_frame
    tf_h.word_wrap = True
    p_h = tf_h.paragraphs[0]
    p_h.text = "MINI PROJECT ABSTRACT PRESENTATION"
    p_h.font.name = "Arial"
    p_h.font.size = Pt(12)
    p_h.font.bold = True
    p_h.font.color.rgb = RGBColor(140, 178, 166) # #8CB2A6 (Mint)
    
    # TITLE (DayTone)
    title_box = slide1.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(8), Inches(1.0))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "DayTone"
    p_t.font.name = "Georgia"
    p_t.font.size = Pt(54)
    p_t.font.bold = True
    p_t.font.color.rgb = RGBColor(255, 255, 255) # White
    
    # SUBTITLE
    subtitle_box = slide1.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(9.5), Inches(0.5))
    tf_s = subtitle_box.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.text = "AI-Powered Mental Wellness Tracker — NLP, ML & Real-Time Analytics"
    p_s.font.name = "Arial"
    p_s.font.size = Pt(14)
    p_s.font.italic = True
    p_s.font.color.rgb = RGBColor(140, 178, 166) # Mint
    
    # TOP RIGHT AUTHOR BOX
    author_bg = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(0.4), Inches(2.5), Inches(1.8))
    author_bg.fill.solid()
    author_bg.fill.fore_color.rgb = RGBColor(28, 59, 54) # #1C3B36 (Translucent Box)
    author_bg.line.color.rgb = RGBColor(28, 59, 54)
    
    tf_a = author_bg.text_frame
    tf_a.word_wrap = True
    tf_a.margin_top = Inches(0.2)
    
    p_a1 = tf_a.paragraphs[0]
    p_a1.text = "S A Sanush"
    p_a1.alignment = PP_ALIGN.CENTER
    p_a1.font.name = "Arial"
    p_a1.font.size = Pt(18)
    p_a1.font.bold = True
    p_a1.font.color.rgb = RGBColor(255, 255, 255)
    
    p_a2 = tf_a.add_paragraph()
    p_a2.text = "S3 MCA"
    p_a2.alignment = PP_ALIGN.CENTER
    p_a2.font.name = "Arial"
    p_a2.font.size = Pt(14)
    p_a2.font.color.rgb = RGBColor(140, 178, 166)
    
    p_a3 = tf_a.add_paragraph()
    p_a3.text = "01/07/2026"
    p_a3.alignment = PP_ALIGN.CENTER
    p_a3.font.name = "Arial"
    p_a3.font.size = Pt(12)
    p_a3.font.color.rgb = RGBColor(140, 178, 166)
    
    # COLUMN 1: About the Project
    col1_left = Inches(0.6)
    col_width = Inches(3.8)
    col_top = Inches(2.6)
    col_height = Inches(4.2)
    
    # Column 1 Header with Icon Circle
    icon1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, col1_left, col_top, Inches(0.5), Inches(0.5))
    icon1.fill.solid()
    icon1.fill.fore_color.rgb = RGBColor(132, 163, 178) # #84A3B2 (Slate Blue)
    icon1.line.color.rgb = RGBColor(132, 163, 178)
    
    # Lightbulb/Icon simple character representation (using text inside oval)
    p_i1 = icon1.text_frame.paragraphs[0]
    p_i1.text = "💡"
    p_i1.alignment = PP_ALIGN.CENTER
    p_i1.font.size = Pt(16)
    
    col1_title = slide1.shapes.add_textbox(col1_left + Inches(0.6), col_top - Inches(0.05), Inches(3.2), Inches(0.6))
    p_c1_t = col1_title.text_frame.paragraphs[0]
    p_c1_t.text = "About the Project"
    p_c1_t.font.name = "Georgia"
    p_c1_t.font.size = Pt(20)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = RGBColor(255, 255, 255)
    
    # Column 1 Content
    col1_content = slide1.shapes.add_textbox(col1_left, col_top + Inches(0.6), col_width, col_height)
    tf_c1 = col1_content.text_frame
    tf_c1.word_wrap = True
    bullets1 = [
        "Flask app with modular blueprints: auth, mood, admin",
        "Daily logs: mood, sleep, stress, activity, social interaction",
        "NLP sentiment scoring — VADER + optional BERT wrapper",
        "ML burnout prediction — Random Forest, ~99% test accuracy",
        "Dashboards: Chart.js, D3.js heatmap, Three.js 3D Mood Orb",
        "PDF reports, streamed CSV export, personal goal tracking"
    ]
    for idx, b_text in enumerate(bullets1):
        p = tf_c1.paragraphs[0] if idx == 0 else tf_c1.add_paragraph()
        p.text = "▪  " + b_text
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_after = Pt(8)
        
    # COLUMN 2: Problem Statement
    col2_left = Inches(4.7)
    
    icon2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, col2_left, col_top, Inches(0.5), Inches(0.5))
    icon2.fill.solid()
    icon2.fill.fore_color.rgb = RGBColor(96, 148, 136) # #609488 (Teal Mint)
    icon2.line.color.rgb = RGBColor(96, 148, 136)
    p_i2 = icon2.text_frame.paragraphs[0]
    p_i2.text = "⚠️"
    p_i2.alignment = PP_ALIGN.CENTER
    p_i2.font.size = Pt(14)
    
    col2_title = slide1.shapes.add_textbox(col2_left + Inches(0.6), col_top - Inches(0.05), Inches(3.2), Inches(0.6))
    p_c2_t = col2_title.text_frame.paragraphs[0]
    p_c2_t.text = "Problem Statement"
    p_c2_t.font.name = "Georgia"
    p_c2_t.font.size = Pt(20)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = RGBColor(255, 255, 255)
    
    col2_content = slide1.shapes.add_textbox(col2_left, col_top + Inches(0.6), col_width, col_height)
    tf_c2 = col2_content.text_frame
    tf_c2.word_wrap = True
    bullets2 = [
        "No early-warning system exists for burnout risk",
        "Manual journaling gives no analytical insight",
        "Mood, sleep & stress patterns go unnoticed over time",
        "Wellness tools are fragmented across separate apps",
        "Little to no personalized, data-driven self-care guidance",
        "Risk labels alone offer no explanation or support"
    ]
    for idx, b_text in enumerate(bullets2):
        p = tf_c2.paragraphs[0] if idx == 0 else tf_c2.add_paragraph()
        p.text = "▪  " + b_text
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_after = Pt(8)
        
    # COLUMN 3: Core Modules
    col3_left = Inches(8.8)
    
    icon3 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, col3_left, col_top, Inches(0.5), Inches(0.5))
    icon3.fill.solid()
    icon3.fill.fore_color.rgb = RGBColor(213, 171, 123) # #D5AB7B (Gold)
    icon3.line.color.rgb = RGBColor(213, 171, 123)
    p_i3 = icon3.text_frame.paragraphs[0]
    p_i3.text = "🥞"
    p_i3.alignment = PP_ALIGN.CENTER
    p_i3.font.size = Pt(14)
    
    col3_title = slide1.shapes.add_textbox(col3_left + Inches(0.6), col_top - Inches(0.05), Inches(3.2), Inches(0.6))
    p_c3_t = col3_title.text_frame.paragraphs[0]
    p_c3_t.text = "Core Modules"
    p_c3_t.font.name = "Georgia"
    p_c3_t.font.size = Pt(20)
    p_c3_t.font.bold = True
    p_c3_t.font.color.rgb = RGBColor(255, 255, 255)
    
    col3_content = slide1.shapes.add_textbox(col3_left, col_top + Inches(0.6), col_width, col_height)
    tf_c3 = col3_content.text_frame
    tf_c3.word_wrap = True
    bullets3 = [
        "Auth — scrypt hashing, admin code, rate-limited login",
        "Mood engine — 12-feature rolling analytics per log",
        "ML pipeline — RF vs Decision Tree vs Logistic Regression",
        "NLP module — VADER sentiment, BERT fallback option",
        "Suggestion engine — rule-based daily 'Do' / 'Don't' tips",
        "Admin module — platform stats & cohort accuracy"
    ]
    for idx, b_text in enumerate(bullets3):
        p = tf_c3.paragraphs[0] if idx == 0 else tf_c3.add_paragraph()
        p.text = "▪  " + b_text
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_after = Pt(8)
        
    # Footer Slide 1
    footer1 = slide1.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(6), Inches(0.4))
    footer1.text_frame.paragraphs[0].text = "DayTone  |  Mini Project Abstract"
    footer1.text_frame.paragraphs[0].font.size = Pt(10)
    footer1.text_frame.paragraphs[0].font.color.rgb = RGBColor(140, 178, 166)
    
    page1 = slide1.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.2), Inches(0.4))
    page1.text_frame.paragraphs[0].text = "01 / 02"
    page1.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    page1.text_frame.paragraphs[0].font.size = Pt(10)
    page1.text_frame.paragraphs[0].font.color.rgb = RGBColor(140, 178, 166)
    
    
    # ---------------------------------------------------------
    # SLIDE 2: Light Theme (Off-White)
    # ---------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(243, 246, 245) # #F3F6F5 (Light Mint Grey)
    
    # Title Slide 2
    title_box2 = slide2.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    tf_t2 = title_box2.text_frame
    tf_t2.word_wrap = True
    p_t2 = tf_t2.paragraphs[0]
    p_t2.text = "Existing System, Security & Technology Stack"
    p_t2.font.name = "Georgia"
    p_t2.font.size = Pt(28) # Reduced to prevent line wrap and overlap
    p_t2.font.bold = True
    p_t2.font.color.rgb = RGBColor(27, 42, 39) # #1B2A27 (Charcoal Green)
    
    # Subtitle Slide 2
    sub2 = slide2.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(10), Inches(0.4))
    p_sub2 = sub2.text_frame.paragraphs[0]
    p_sub2.text = "DayTone  |  Mini Project Abstract Presentation"
    p_sub2.font.name = "Arial"
    p_sub2.font.size = Pt(12)
    p_sub2.font.italic = True
    p_sub2.font.color.rgb = RGBColor(96, 148, 136) # Teal Mint
    
    # Columns positioning
    col_y = Inches(2.1)
    col_h = Inches(4.4)
    
    # Col 1: Existing System
    c1_left = Inches(0.6)
    icon_s2_1 = slide2.shapes.add_shape(MSO_SHAPE.OVAL, c1_left, col_y, Inches(0.5), Inches(0.5))
    icon_s2_1.fill.solid()
    icon_s2_1.fill.fore_color.rgb = RGBColor(132, 163, 178) # Slate Blue
    icon_s2_1.line.color.rgb = RGBColor(132, 163, 178)
    icon_s2_1.text_frame.paragraphs[0].text = "💾"
    icon_s2_1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon_s2_1.text_frame.paragraphs[0].font.size = Pt(14)
    
    c1_title = slide2.shapes.add_textbox(c1_left + Inches(0.6), col_y - Inches(0.05), Inches(3.2), Inches(0.6))
    p_c1_t2 = c1_title.text_frame.paragraphs[0]
    p_c1_t2.text = "Existing System"
    p_c1_t2.font.name = "Georgia"
    p_c1_t2.font.size = Pt(20)
    p_c1_t2.font.bold = True
    p_c1_t2.font.color.rgb = RGBColor(27, 42, 39)
    
    c1_content = slide2.shapes.add_textbox(c1_left, col_y + Inches(0.6), col_width, col_h)
    tf_c1_s2 = c1_content.text_frame
    tf_c1_s2.word_wrap = True
    bullets_s2_1 = [
        "Plain journaling / diary apps with no analysis",
        "Generic mood trackers — manual entry only",
        "No predictive burnout-risk classification",
        "No NLP-based sentiment evaluation of notes",
        "Disconnected tools across sleep, stress & mood",
        "Limited or no personalized recommendations"
    ]
    for idx, b_text in enumerate(bullets_s2_1):
        p = tf_c1_s2.paragraphs[0] if idx == 0 else tf_c1_s2.add_paragraph()
        p.text = "▪  " + b_text
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(8)
        
    # Col 2: Roles & Security
    c2_left = Inches(4.7)
    icon_s2_2 = slide2.shapes.add_shape(MSO_SHAPE.OVAL, c2_left, col_y, Inches(0.5), Inches(0.5))
    icon_s2_2.fill.solid()
    icon_s2_2.fill.fore_color.rgb = RGBColor(96, 148, 136) # Teal Mint
    icon_s2_2.line.color.rgb = RGBColor(96, 148, 136)
    icon_s2_2.text_frame.paragraphs[0].text = "🛡️"
    icon_s2_2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon_s2_2.text_frame.paragraphs[0].font.size = Pt(14)
    
    c2_title = slide2.shapes.add_textbox(c2_left + Inches(0.6), col_y - Inches(0.05), Inches(3.2), Inches(0.6))
    p_c2_t2 = c2_title.text_frame.paragraphs[0]
    p_c2_t2.text = "Roles & Security"
    p_c2_t2.font.name = "Georgia"
    p_c2_t2.font.size = Pt(20)
    p_c2_t2.font.bold = True
    p_c2_t2.font.color.rgb = RGBColor(27, 42, 39)
    
    c2_content = slide2.shapes.add_textbox(c2_left, col_y + Inches(0.6), col_width, col_h)
    tf_c2_s2 = c2_content.text_frame
    tf_c2_s2.word_wrap = True
    bullets_s2_2 = [
        "Two roles: Regular User and Administrator",
        "CSRF protection (Flask-WTF) & scrypt password hashing",
        "CSP headers (Flask-Talisman) + login rate limiting",
        "Session hardening: HttpOnly, SameSite=Lax, 8-hr lifetime",
        "Safe-redirect validation, 2 MB request size cap",
        "Production blocks SQLite — requires PostgreSQL"
    ]
    for idx, b_text in enumerate(bullets_s2_2):
        p = tf_c2_s2.paragraphs[0] if idx == 0 else tf_c2_s2.add_paragraph()
        p.text = "▪  " + b_text
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(8)
        
    # Col 3: Technology Used (Cards Style)
    c3_left = Inches(8.8)
    icon_s2_3 = slide2.shapes.add_shape(MSO_SHAPE.OVAL, c3_left, col_y, Inches(0.5), Inches(0.5))
    icon_s2_3.fill.solid()
    icon_s2_3.fill.fore_color.rgb = RGBColor(213, 171, 123) # Gold
    icon_s2_3.line.color.rgb = RGBColor(213, 171, 123)
    icon_s2_3.text_frame.paragraphs[0].text = "💼"
    icon_s2_3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon_s2_3.text_frame.paragraphs[0].font.size = Pt(14)
    
    c3_title = slide2.shapes.add_textbox(c3_left + Inches(0.6), col_y - Inches(0.05), Inches(3.2), Inches(0.6))
    p_c3_t2 = c3_title.text_frame.paragraphs[0]
    p_c3_t2.text = "Technology Used"
    p_c3_t2.font.name = "Georgia"
    p_c3_t2.font.size = Pt(20)
    p_c3_t2.font.bold = True
    p_c3_t2.font.color.rgb = RGBColor(27, 42, 39)
    
    tech_cards = [
        ("Backend", "Flask, SQLAlchemy, Alembic migrations"),
        ("ML / NLP", "Scikit-Learn (RF/DT/LogReg), NLTK VADER, BERT"),
        ("Visualization", "Chart.js, D3.js, Three.js"),
        ("Security", "Flask-Talisman, Flask-Limiter, Fernet"),
        ("DevOps", "Gunicorn, GitHub Actions CI, pytest (25 tests)")
    ]
    
    card_y = col_y + Inches(0.6)
    card_h = Inches(0.72)
    card_gap = Inches(0.12)
    
    for c_title, c_desc in tech_cards:
        # draw rounded box
        card_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c3_left, card_y, col_width, card_h)
        card_box.fill.solid()
        card_box.fill.fore_color.rgb = RGBColor(235, 240, 238) # #EBF0EE (Light Card Grey)
        card_box.line.color.rgb = RGBColor(220, 226, 224)
        
        tf_card = card_box.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.15)
        tf_card.margin_right = Inches(0.15)
        tf_card.margin_top = Inches(0.1)
        tf_card.margin_bottom = Inches(0.1)
        
        p1 = tf_card.paragraphs[0]
        p1.text = c_title
        p1.font.name = "Arial"
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(27, 42, 39)
        
        p2 = tf_card.add_paragraph()
        p2.text = c_desc
        p2.font.name = "Arial"
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(80, 80, 80)
        
        card_y += card_h + card_gap
        
    # Footer Slide 2
    footer2 = slide2.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(6), Inches(0.4))
    footer2.text_frame.paragraphs[0].text = "DayTone  |  Mini Project Abstract"
    footer2.text_frame.paragraphs[0].font.size = Pt(10)
    footer2.text_frame.paragraphs[0].font.color.rgb = RGBColor(96, 148, 136)
    
    page2 = slide2.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.2), Inches(0.4))
    page2.text_frame.paragraphs[0].text = "02 / 02"
    page2.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    page2.text_frame.paragraphs[0].font.size = Pt(10)
    page2.text_frame.paragraphs[0].font.color.rgb = RGBColor(96, 148, 136)
    
    # Save presentation
    output_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "DayTone_Abstract_Presentation.pptx"))
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
